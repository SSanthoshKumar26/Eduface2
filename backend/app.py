import os
import subprocess
import sys

def install_and_import(package):
    try:
        __import__(package)
    except ImportError:
        print(f"[*] Installing {package} for premium features (this may take a minute)...")
        try:
            subprocess.check_call([sys.executable, "-m", "pip", "install", package, "--quiet", "--no-input"])
        except:
            pass

# Auto-install PDF/DOCX exporters if missing
try: install_and_import('fpdf2')
except: pass

try: install_and_import('python-docx')
except: pass
import sys
import io
import threading

# Fix terminal encoding issues on Windows
try:
    if sys.stdout.encoding != 'utf-8':
        import codecs
        sys.stdout = codecs.getwriter('utf-8')(sys.stdout.buffer, 'strict')
        sys.stderr = codecs.getwriter('utf-8')(sys.stderr.buffer, 'strict')
except:
    pass

# Pre-configure ffmpeg path for pydub to avoid RuntimeWarning
# We must do this before importing video_processor
backend_dir = os.path.dirname(os.path.abspath(__file__))
ffmpeg_dir = os.path.join(backend_dir, 'Wav2Lip')
if os.path.exists(os.path.join(ffmpeg_dir, 'ffmpeg.exe')):
    if ffmpeg_dir not in os.environ["PATH"]:
        os.environ["PATH"] = ffmpeg_dir + os.pathsep + os.environ["PATH"]
import re
import requests
import uuid
import json
import time
from flask import Flask, request, jsonify, send_file
from flask_cors import CORS
from dotenv import load_dotenv
import requests
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
from io import BytesIO
from werkzeug.utils import secure_filename
from openai import OpenAI
from groq import Groq

# Video generation imports
try:
    from video_processor.pipeline import VideoPipeline
    VIDEO_GENERATION_ENABLED = True
except ImportError:
    VIDEO_GENERATION_ENABLED = False
    print("[WARNING] Video generation modules not found. Video features will be disabled.")

load_dotenv()

app = Flask(__name__)
CORS(app)

# ============ CONFIGURATION ============
# AI Model Configuration (Ollama)
OLLAMA_BASE_URL = os.getenv("OLLAMA_BASE_URL", "http://localhost:11434")
AI_MODEL = os.getenv("OLLAMA_MODEL", "tinyllama")

# Image APIs
UNSPLASH_API_KEY = os.getenv("UNSPLASH_API_KEY", "").strip()
PEXELS_API_KEY = os.getenv("PEXELS_API_KEY", "").strip()
PIXABAY_API_KEY = os.getenv("PIXABAY_API_KEY", "").strip()

# File upload configuration
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
UPLOAD_FOLDER = os.path.join(BASE_DIR, 'uploads')
MAX_CONTENT_LENGTH = 100 * 1024 * 1024  # 100MB
ALLOWED_PPT_EXTENSIONS = {'ppt', 'pptx'}
ALLOWED_FACE_EXTENSIONS = {'jpg', 'jpeg', 'png', 'mp4', 'mov', 'webm'}


app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['MAX_CONTENT_LENGTH'] = MAX_CONTENT_LENGTH

# Create upload directories
PPT_UPLOAD_DIR = os.path.join(UPLOAD_FOLDER, 'ppts')
FACE_UPLOAD_DIR = os.path.join(UPLOAD_FOLDER, 'faces')
OUTPUT_DIR = os.path.join(UPLOAD_FOLDER, 'outputs')
SHARED_CHATS_DIR = os.path.join(UPLOAD_FOLDER, 'shared_chats')

os.makedirs(PPT_UPLOAD_DIR, exist_ok=True)
os.makedirs(FACE_UPLOAD_DIR, exist_ok=True)
os.makedirs(OUTPUT_DIR, exist_ok=True)
os.makedirs(SHARED_CHATS_DIR, exist_ok=True)

# Initialize video pipeline if available
video_pipeline = None
if VIDEO_GENERATION_ENABLED:
    try:
        video_pipeline = VideoPipeline(output_dir=OUTPUT_DIR)
        print("[SUCCESS] Video generation pipeline initialized")
    except Exception as e:
        print(f"[WARNING] Could not initialize video pipeline: {e}")
        VIDEO_GENERATION_ENABLED = False

# ============ MODE CONFIGURATIONS ============
MODE_PROMPTS = {
    "Quick Response": {
        "instructions": "Create a concise presentation with key points only. Keep it brief and to the point.",
        "slide_count": 3,
        "max_bullets": 3,
        "font_size_body": 13
    },
    "Creative": {
        "instructions": "Create an engaging and creative presentation with storytelling elements. Use vivid descriptions and make it interesting.",
        "slide_count": 5,
        "max_bullets": 4,
        "font_size_body": 12
    },
    "Detailed": {
        "instructions": "Create a comprehensive presentation with detailed explanations, examples, and in-depth analysis. Cover all aspects thoroughly.",
        "slide_count": 7,
        "max_bullets": 5,
        "font_size_body": 11
    }
}

# ============ THEMES ============
THEMES = {
    "modern_blue": {
        "name": "Modern Blue",
        "colors": {
            "background": (245, 245, 255),
            "title_slide_bg": (19, 71, 130),
            "title": (255, 255, 255),
            "subtitle": (220, 230, 250),
            "slide_title": (19, 71, 130),
            "h2": (30, 90, 150),
            "h3": (70, 120, 170),
            "text": (45, 45, 45),
            "bullet": (60, 60, 60),
            "accent": (19, 71, 130),
            "accent_light": (200, 220, 255),
            "accent_dark": (10, 40, 80)
        },
        "fonts": {"title": "Calibri", "heading": "Calibri", "text": "Calibri"}
    },
    "modern_purple": {
        "name": "Elegant Purple",
        "colors": {
            "background": (240, 235, 250),
            "title_slide_bg": (102, 51, 153),
            "title": (255, 255, 255),
            "subtitle": (220, 200, 240),
            "slide_title": (102, 51, 153),
            "h2": (120, 70, 170),
            "h3": (150, 120, 190),
            "text": (65, 65, 65),
            "bullet": (75, 75, 75),
            "accent": (102, 51, 153),
            "accent_light": (230, 210, 250),
            "accent_dark": (70, 30, 110)
        },
        "fonts": {"title": "Times New Roman", "heading": "Times New Roman", "text": "Times New Roman"}
    },
    "modern_green": {
        "name": "Creative Green",
        "colors": {
            "background": (235, 245, 235),
            "title_slide_bg": (34, 102, 36),
            "title": (255, 255, 255),
            "subtitle": (200, 230, 200),
            "slide_title": (34, 102, 36),
            "h2": (60, 130, 65),
            "h3": (100, 150, 105),
            "text": (50, 50, 50),
            "bullet": (60, 60, 60),
            "accent": (34, 102, 36),
            "accent_light": (200, 240, 200),
            "accent_dark": (20, 70, 25)
        },
        "fonts": {"title": "Segoe UI", "heading": "Segoe UI", "text": "Segoe UI"}
    },
    "modern_sunset": {
        "name": "Vibrant Orange",
        "colors": {
            "background": (255, 245, 235),
            "title_slide_bg": (204, 85, 0),
            "title": (255, 255, 255),
            "subtitle": (255, 220, 180),
            "slide_title": (204, 85, 0),
            "h2": (220, 110, 30),
            "h3": (240, 150, 80),
            "text": (60, 40, 20),
            "bullet": (80, 60, 40),
            "accent": (204, 85, 0),
            "accent_light": (255, 220, 180),
            "accent_dark": (150, 60, 0)
        },
        "fonts": {"title": "Helvetica", "heading": "Helvetica", "text": "Helvetica"}
    },
    "minimal_dark": {
        "name": "Minimalist Black",
        "colors": {
            "background": (20, 20, 20),
            "title_slide_bg": (0, 0, 0),
            "title": (255, 255, 255),
            "subtitle": (200, 200, 200),
            "slide_title": (255, 255, 255),
            "h2": (220, 220, 220),
            "h3": (180, 180, 180),
            "text": (180, 180, 180),
            "bullet": (160, 160, 160),
            "accent": (255, 255, 255),
            "accent_light": (60, 60, 60),
            "accent_dark": (100, 100, 100)
        },
        "fonts": {"title": "Arial Black", "heading": "Arial Black", "text": "Arial"}
    },
    "minimal_light": {
        "name": "Minimalist White",
        "colors": {
            "background": (255, 255, 255),
            "title_slide_bg": (240, 240, 240),
            "title": (30, 30, 30),
            "subtitle": (100, 100, 100),
            "slide_title": (0, 0, 0),
            "h2": (50, 50, 50),
            "h3": (100, 100, 100),
            "text": (60, 60, 60),
            "bullet": (80, 80, 80),
            "accent": (150, 150, 150),
            "accent_light": (245, 245, 245),
            "accent_dark": (50, 50, 50)
        },
        "fonts": {"title": "Arial", "heading": "Arial", "text": "Arial"}
    },
    "corporate": {
        "name": "Corporate Gray",
        "colors": {
            "background": (250, 250, 250),
            "title_slide_bg": (70, 70, 70),
            "title": (255, 255, 255),
            "subtitle": (220, 220, 220),
            "slide_title": (50, 50, 50),
            "h2": (80, 80, 80),
            "h3": (120, 120, 120),
            "text": (60, 60, 60),
            "bullet": (70, 70, 70),
            "accent": (150, 150, 150),
            "accent_light": (230, 230, 230),
            "accent_dark": (40, 40, 40)
        },
        "fonts": {"title": "Arial", "heading": "Arial", "text": "Arial"}
    },
    "creative": {
        "name": "Cool Teal",
        "colors": {
            "background": (230, 248, 248),
            "title_slide_bg": (0, 102, 102),
            "title": (255, 255, 255),
            "subtitle": (180, 220, 220),
            "slide_title": (0, 102, 102),
            "h2": (20, 130, 130),
            "h3": (60, 150, 150),
            "text": (40, 70, 70),
            "bullet": (50, 80, 80),
            "accent": (0, 102, 102),
            "accent_light": (200, 240, 240),
            "accent_dark": (0, 70, 70)
        },
        "fonts": {"title": "Verdana", "heading": "Verdana", "text": "Verdana"}
    },
}

# ============ SLIDE DIMENSIONS ============
SLIDE_WIDTH = 10.0
SLIDE_HEIGHT = 7.5
MARGIN_LEFT = 0.5
MARGIN_RIGHT = 0.5
MARGIN_TOP = 0.4
MARGIN_BOTTOM = 0.5
CONTENT_WIDTH = SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT
TITLE_AREA_HEIGHT = 1.0
CONTENT_AREA_TOP = MARGIN_TOP + TITLE_AREA_HEIGHT + 0.2
AVAILABLE_CONTENT_HEIGHT = SLIDE_HEIGHT - MARGIN_TOP - MARGIN_BOTTOM - TITLE_AREA_HEIGHT - 0.5
TEXT_WIDTH_WITH_IMAGE = 4.8
TEXT_WIDTH_WITHOUT_IMAGE = CONTENT_WIDTH
IMAGE_START_LEFT = 5.3
IMAGE_GAP = 0.2

# ============ UTILITY FUNCTIONS ============
def check_ollama_connection():
    """Verify Ollama is running and accessible"""
    try:
        response = requests.get(f"{OLLAMA_BASE_URL}/api/tags", timeout=2)
        if response.status_code == 200:
            models = response.json().get('models', [])
            model_names = [m.get('name', '').split(':')[0] for m in models]
            return True, model_names
        return False, []
    except Exception as e:
        return False, []

def allowed_file(filename, allowed_extensions):
    """Check if file has allowed extension"""
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in allowed_extensions

def startup_diagnostics():
    """Run startup checks and print status"""
    print("\n" + "="*60)
    print("🚀 FACEPREP - AI PRESENTATION & VIDEO GENERATOR")
    print("="*60)
    
    ollama_ok, models = check_ollama_connection()
    
    if ollama_ok:
        print(f"[SUCCESS] Ollama is running ({OLLAMA_BASE_URL})")
        if models:
            print(f"[*] Available models: {models}")
            if AI_MODEL in models:
                print(f"[SUCCESS] {AI_MODEL} model found!")
            else:
                print(f"[WARNING] {AI_MODEL} not found. Run: ollama pull {AI_MODEL}")
        else:
            print(f"[WARNING] No models found. Run: ollama pull {AI_MODEL}")
    else:
        print(f"[ERROR] Ollama not running at {OLLAMA_BASE_URL}")
        print("   Start with: ollama serve")
    
    print("\n[STATUS] Image API Status:")
    print(f"  {'[OK]' if UNSPLASH_API_KEY else '[MISSING]'} Unsplash API")
    print(f"  {'[OK]' if PEXELS_API_KEY else '[MISSING]'} Pexels API")
    print(f"  {'[OK]' if PIXABAY_API_KEY else '[MISSING]'} Pixabay API")
    
    print("\n[STATUS] Video Generation Status:")
    if VIDEO_GENERATION_ENABLED:
        print("  [SUCCESS] Video generation modules loaded")
        print("  [SUCCESS] Video pipeline initialized")
    else:
        print("  [MISSING] Video generation disabled (modules not found)")
    
    print("\n[SUCCESS] Server ready!")
    print(f"[SUCCESS] Running on: http://localhost:5000")
    print("="*60 + "\n")

startup_diagnostics()

# ============ FONT SIZE CALCULATOR ============
def calculate_dynamic_font_sizes(slide_sections, mode, has_image):
    """Dynamically calculate font sizes based on content density"""
    bullet_count = sum(1 for s in slide_sections if s['type'] == 'bullet')
    text_count = sum(1 for s in slide_sections if s['type'] == 'text')
    h3_count = sum(1 for s in slide_sections if s['type'] == 'h3')
    
    total_items = bullet_count + text_count + h3_count
    
    base_body = MODE_PROMPTS[mode].get("font_size_body", 12)
    
    if total_items <= 2:
        body_size = base_body + 1
        h3_size = 16
        spacing = 8
    elif total_items <= 4:
        body_size = base_body
        h3_size = 15
        spacing = 6
    elif total_items <= 6:
        body_size = max(base_body - 0.5, 11)
        h3_size = 14
        spacing = 5
    elif total_items <= 8:
        body_size = max(base_body - 1, 10)
        h3_size = 13
        spacing = 4
    elif total_items <= 12:
        body_size = max(base_body - 1.5, 9)
        h3_size = 12
        spacing = 3
    else:
        body_size = max(base_body - 2, 8)
        h3_size = 11
        spacing = 2
    
    if has_image:
        body_size = max(body_size - 0.5, 8)
        h3_size = max(h3_size - 1, 10)
        spacing = max(spacing - 1, 1)
    
    return {
        'body': body_size,
        'h3': h3_size,
        'spacing': spacing,
        'line_spacing': 1.1 if total_items > 8 else 1.15
    }

# ============ AI GENERATION ============
def generate_with_ai(prompt: str, mode: str = "Creative", slide_count: int = 5) -> str:
    """Generate content trying Groq first, then fallback to local AI model via Ollama"""
    mode_config = MODE_PROMPTS.get(mode, MODE_PROMPTS["Creative"])
    mode_instruction = mode_config.get("instructions", "")
    
    enhanced_prompt = f"""{mode_instruction}

User Topic/Request: {prompt}

Please structure your response ONLY with markdown:
- Use # for the main topic (only once at the very start)
- Use ## for each major section (each ## becomes one slide)
- Create EXACTLY {slide_count} slides (exactly {slide_count} ## sections)
- Use ### for subsections within a slide
- Use * or - for bullet points (keep to {mode_config.get('max_bullets', 4)} per section)
- Keep points clear, concise, and well-organized
- Include practical examples where relevant

CRITICAL INSTRUCTIONS:
- Do NOT mention "Slide 1", "Slide 2", etc in the content
- Do NOT use dashes or separators like "-----" or "============"
- Only use proper markdown headers (# ## ###)
- Each ## section MUST be a separate slide
- Focus on clear, structured markdown ONLY
- Generate comprehensive content for ALL {slide_count} slides"""

    # Try Groq API first
    groq_api_key = os.getenv("GROQ_API_KEY", "").strip()
    if groq_api_key:
        try:
            print(f"🤖 Generating with Eduface AI Persona (Mode: {mode})...")
            client = Groq(api_key=groq_api_key)
            
            system_prompt = """You are Eduface AI — an advanced conversational learning and content creation assistant.
You operate as a REAL-TIME INTERACTIVE CHAT SYSTEM with a persistent working document.

CORE MODES:
1. CHAT MODE (conversation with user)
2. DOCUMENT MODE (live structured content)

Every response MUST have 2 parts:
1. [CHAT] - A short, natural, conversational reply acknowledging the user (1-2 lines max).
2. [DOCUMENT] - The full, updated, structured markdown content.

RULES:
- Maintain a persistent document. All updates modify the existing structure.
- [DOCUMENT] SECTION MUST ONLY CONTAIN THE FINAL TOPIC CONTENT.
- DO NOT include ANY conversational meta-talk in the [DOCUMENT] segment (e.g., Avoid: 'Here is your update...', 'Let me know if you need more...').
- Put ALL conversational pleasantries, explanations, and follow-up questions exclusively in the [CHAT] segment.
- If input is voice, the [CHAT] part should be warmer, but the [DOCUMENT] part must remain a clean, professional written asset.
- Use structured markdown (Headings, Bullet points) in [DOCUMENT]."""

            response = client.chat.completions.create(
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": f"Mode: {mode}\nUser Input: {prompt}"}
                ],
                model="llama-3.3-70b-versatile",
                temperature=0.7,
            )
            
            if response.choices and len(response.choices) > 0:
                generated_text = response.choices[0].message.content.strip()
                print(f"✅ Generation complete via Groq ({len(generated_text)} chars)\n")
                return generated_text
        except Exception as e:
            print(f"[WARNING] Groq API failed: {str(e)}. Falling back to Ollama...")
    else:
        print("[INFO] No GROQ_API_KEY found. Falling back to Ollama...")
    
    # Fallback to Ollama
    try:
        url = f"{OLLAMA_BASE_URL}/api/generate"
        payload = {
            "model": AI_MODEL,
            "prompt": enhanced_prompt,
            "stream": False,
            "temperature": 0.7 if mode == "Creative" else (0.5 if mode == "Quick Response" else 0.6),
        }
        
        print(f"🤖 Generating [{mode}] ({slide_count} slides) with local {AI_MODEL}...")
        
        # Increased timeout to 300s to allow the 1B model to finish without throwing an error
        response = requests.post(url, json=payload, timeout=300)
        
        if response.status_code == 200:
            result = response.json()
            generated_text = result.get('response', '').strip()
            print(f"✅ Generation complete ({len(generated_text)} chars)\n")
            return generated_text
        else:
            error_msg = f"Ollama error: HTTP {response.status_code}"
            print(f"❌ {error_msg}")
            return f"[Error]: {error_msg}"
            
    except requests.exceptions.ConnectionError:
        error_msg = f"Cannot connect to Ollama at {OLLAMA_BASE_URL}."
        print(f"❌ {error_msg}")
        return f"[Error]: {error_msg}"
    
    except requests.exceptions.Timeout:
        error_msg = f"Generation timed out (300s). Try a simpler prompt or Quick Response mode."
        print(f"❌ {error_msg}")
        return f"[Error]: {error_msg}"
    
    except Exception as e:
        error_msg = str(e)
        print(f"❌ AI error: {error_msg}")
        return f"[Error]: {error_msg}"
# ============ TEXT PROCESSING ============
def clean_markdown(text):
    """Remove markdown formatting"""
    if not text:
        return ""
    
    text = str(text)
    text = re.sub(r'^[\~\-\*]{3,}$', '', text, flags=re.MULTILINE)
    text = re.sub(r'\*\*(.*?)\*\*', r'\1', text)
    text = re.sub(r'__(.*?)__', r'\1', text)
    text = re.sub(r'\*(.*?)\*', r'\1', text)
    text = re.sub(r'_(.*?)_', r'\1', text)
    text = re.sub(r'`(.*?)`', r'\1', text)
    text = re.sub(r'```(.*?)```', r'\1', text, flags=re.DOTALL)
    text = re.sub(r'#+\s*', '', text)
    text = re.sub(r'^[\*\-\•]\s*', '', text, flags=re.MULTILINE)
    text = re.sub(r'\[([^\]]+)\]\(([^\)]+)\)', r'\1', text)
    text = re.sub(r'\[\s*\]', '', text)
    text = re.sub(r'\(\s*\)', '', text)
    text = re.sub(r'\s+', ' ', text)
    text = re.sub(r'\n+', '\n', text)
    
    return text.strip()

def truncate_text(text, max_length=250):
    """Truncate text to prevent overflow"""
    text = text.strip()
    if len(text) > max_length:
        truncated = text[:max_length]
        last_space = truncated.rfind(' ')
        if last_space > max_length * 0.7:
            return text[:last_space] + "..."
        return truncated + "..."
    return text

def parse_markdown_content(text):
    """Parse markdown content into sections"""
    lines = text.split('\n')
    sections = []
    
    i = 0
    while i < len(lines):
        line = lines[i]
        stripped = line.strip()
        
        if not stripped or len(stripped) < 2:
            i += 1
            continue
        
        if re.match(r'^[\~\-\*]{3,}$', stripped):
            i += 1
            continue
        
        if stripped.startswith('###'):
            content = stripped.lstrip('#').strip()
            content = clean_markdown(content)
            if content and len(content) > 2:
                sections.append({'type': 'h3', 'content': content, 'level': 3})
        
        elif stripped.startswith('##'):
            content = stripped.lstrip('#').strip()
            content = clean_markdown(content)
            if content and len(content) > 2:
                sections.append({'type': 'h2', 'content': content, 'level': 2})
        
        elif stripped.startswith('#'):
            content = stripped.lstrip('#').strip()
            content = clean_markdown(content)
            if content and len(content) > 2:
                sections.append({'type': 'h1', 'content': content, 'level': 1})
        
        elif stripped.startswith(('*', '-', '•')):
            content = stripped.lstrip('*-•').strip()
            content = clean_markdown(content)
            if content and len(content) > 2:
                level = (len(line) - len(line.lstrip())) // 2
                sections.append({'type': 'bullet', 'content': content, 'level': level})
        
        else:
            content = clean_markdown(stripped)
            if len(content) > 5:
                sections.append({'type': 'text', 'content': content, 'level': 0})
        
        i += 1
    
    return sections

def group_into_slides(sections, max_slides=5):
    """Smartly group sections into EXACTLY max_slides without empty slides"""
    valid_sections = [s for s in sections if str(s.get('content', '')).strip()]
    if not valid_sections:
        return []
        
    raw_slides = []
    current_slide = []
    
    for s in valid_sections:
        if s['type'] in ['h1', 'h2']:
            if current_slide:
                raw_slides.append(current_slide)
            current_slide = [s]
        else:
            current_slide.append(s)
            
    if current_slide:
        raw_slides.append(current_slide)
        
    clean_slides = []
    for slide in raw_slides:
        content_items = [x for x in slide if x['type'] in ['text', 'bullet', 'h3']]
        if content_items:
            clean_slides.append(slide)
            
    if not clean_slides:
        clean_slides = [s for s in raw_slides if s]
        
    if not clean_slides:
        return []

    while len(clean_slides) > max_slides:
        min_len = float('inf')
        min_idx = 0
        for i in range(len(clean_slides) - 1):
            cl = len(clean_slides[i]) + len(clean_slides[i+1])
            if cl < min_len:
                min_len = cl
                min_idx = i
        merged = clean_slides[min_idx] + clean_slides[min_idx+1]
        clean_slides[min_idx] = merged
        del clean_slides[min_idx+1]

    while len(clean_slides) < max_slides:
        longest_idx = max(range(len(clean_slides)), key=lambda i: len([x for x in clean_slides[i] if x['type'] in ['text', 'bullet']]))
        longest = clean_slides[longest_idx]
        text_items = [x for x in longest if x['type'] in ['text', 'bullet']]
        
        if len(text_items) < 2:
            break
            
        mid = len(longest) // 2
        while mid > 0 and longest[mid-1]['type'] in ['h1', 'h2', 'h3']:
            mid -= 1
        if mid == 0:
            mid = 1
            
        slide1 = longest[:mid]
        slide2 = longest[mid:]
        
        header = next((x for x in slide1 if x['type'] in ['h1', 'h2']), None)
        if header and not any(x['type'] in ['h1', 'h2'] for x in slide2):
            slide2.insert(0, {'type': header['type'], 'content': f"{header['content']} (Cont.)"})
            
        clean_slides[longest_idx] = slide1
        clean_slides.insert(longest_idx + 1, slide2)
        
    return clean_slides[:max_slides]

# ============ IMAGE FUNCTIONS ============
def fetch_unsplash_image(query):
    if not UNSPLASH_API_KEY:
        return None
    try:
        url = "https://api.unsplash.com/photos/random"
        params = {"query": query, "client_id": UNSPLASH_API_KEY, "orientation": "landscape"}
        response = requests.get(url, params=params, timeout=5)
        if response.status_code == 200:
            return response.json().get('urls', {}).get('regular')
    except:
        pass
    return None

def fetch_pexels_image(query):
    if not PEXELS_API_KEY:
        return None
    try:
        url = "https://api.pexels.com/v1/search"
        headers = {"Authorization": PEXELS_API_KEY}
        params = {"query": query, "per_page": 1}
        response = requests.get(url, headers=headers, params=params, timeout=5)
        if response.status_code == 200:
            photos = response.json().get('photos', [])
            if photos:
                return photos[0].get('src', {}).get('large')
    except:
        pass
    return None

def fetch_pixabay_image(query):
    if not PIXABAY_API_KEY:
        return None
    try:
        url = "https://pixabay.com/api/"
        params = {"key": PIXABAY_API_KEY, "q": query, "image_type": "photo", "per_page": 1, "min_width": 800, "min_height": 600}
        response = requests.get(url, params=params, timeout=5)
        if response.status_code == 200:
            hits = response.json().get('hits', [])
            if hits:
                return hits[0].get('largeImageURL')
    except:
        pass
    return None

def get_image_for_slide(slide_title):
    search_query = re.sub(r'[^a-zA-Z0-9\s]', '', slide_title).strip()[:50]
    if not search_query:
        return None
    
    if UNSPLASH_API_KEY:
        img_url = fetch_unsplash_image(search_query)
        if img_url:
            return img_url
    
    if PEXELS_API_KEY:
        img_url = fetch_pexels_image(search_query)
        if img_url:
            return img_url
    
    if PIXABAY_API_KEY:
        img_url = fetch_pixabay_image(search_query)
        if img_url:
            return img_url
    
    return None

def download_image(url):
    try:
        response = requests.get(url, timeout=10)
        if response.status_code == 200:
            img = Image.open(BytesIO(response.content))
            img.thumbnail((1200, 800), Image.Resampling.LANCZOS)
            img_byte_arr = BytesIO()
            img.save(img_byte_arr, format='JPEG', quality=85)
            img_byte_arr.seek(0)
            return img_byte_arr
    except:
        pass
    return None

# ============ SLIDE BUILDING ============
def add_decorative_shapes(slide, theme, position="top"):
    if position == "top":
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(SLIDE_WIDTH), Inches(0.08))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*theme["colors"]["accent"])
        shape.line.color.rgb = RGBColor(*theme["colors"]["accent"])

def set_background(slide, rgb):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*rgb)

# ============ API ROUTES - CONTENT GENERATION ============
@app.route("/api/generate", methods=["POST"])
def generate_content():
    """Generate content using requested mode and model"""
    data = request.get_json()
    prompt = data.get("prompt", "")
    context = data.get("context", "")
    mode = data.get("mode", "Creative")
    slide_count = data.get("slide_count", 5)
    
    if context:
        prompt = f"Existing document (for context):\n{context}\n\nUser's new request/update: {prompt}"
    
    if mode not in MODE_PROMPTS:
        mode = "Creative"
    
    print(f"\n📝 Generating content in '{mode}' mode for {slide_count} slides...")
    output = generate_with_ai(prompt, mode, slide_count)
    
    return jsonify({
        "output": output,
        "mode": mode,
        "status": "success" if not output.startswith("[Error]") else "error"
    })

@app.route("/api/themes", methods=["GET"])
def get_themes():
    """Get available themes"""
    theme_list = [
        {
            "id": k,
            "name": v["name"],
            "preview_color": '#%02x%02x%02x' % v["colors"]["background"]
        }
        for k, v in THEMES.items()
    ]
    return jsonify({"themes": theme_list})

@app.route("/api/modes", methods=["GET"])
def get_modes():
    """Get available generation modes"""
    modes_list = [
        {
            "id": mode,
            "name": mode,
            "description": config.get("instructions"),
            "slide_count": config.get("slide_count"),
            "max_bullets": config.get("max_bullets")
        }
        for mode, config in MODE_PROMPTS.items()
    ]
    return jsonify({"modes": modes_list})

@app.route("/api/generate-ppt", methods=["POST"])
def generate_ppt():
    """Generate PowerPoint from content"""
    try:
        data = request.get_json()
        content = data.get("content", "").strip()
        title = data.get("title", "Generated Presentation").strip()
        filename = (data.get("filename", "") or "").strip().replace(" ", "_")
        include_images = data.get("include_images", True)
        mode = data.get("mode", "Creative")
        
        if mode not in MODE_PROMPTS:
            mode = "Creative"
        
        if not filename:
            filename = title.replace(" ", "_")
        
        customizations = data.get("customizations", {}) or {}
        theme_id = customizations.get("theme", "modern_blue")
        theme = THEMES.get(theme_id, THEMES["modern_blue"])
        max_slides = int(customizations.get("slide_count", MODE_PROMPTS[mode]["slide_count"]))
        
        if not content:
            return jsonify({"error": "Content is empty"}), 400
        
        sections = parse_markdown_content(content)
        if not sections:
            return jsonify({"error": "Could not parse content"}), 400
        
        slides_content = group_into_slides(sections, max_slides)
        
        prs = Presentation()
        prs.slide_width = Inches(SLIDE_WIDTH)
        prs.slide_height = Inches(SLIDE_HEIGHT)
        
        print("\n" + "="*60)
        print("🎬 CREATING PRESENTATION")
        print("="*60)
        print(f"Title: {title}")
        print(f"Mode: {mode}")
        print(f"Theme: {theme_id}")
        print(f"Total Slides: {len(slides_content)}")
        print("="*60 + "\n")
        
        # Title slide
        title_slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_background(title_slide, theme["colors"]["title_slide_bg"])
        add_decorative_shapes(title_slide, theme, "top")
        
        title_box = title_slide.shapes.add_textbox(
            Inches(MARGIN_LEFT),
            Inches(SLIDE_HEIGHT * 0.25),
            Inches(CONTENT_WIDTH),
            Inches(1.5)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.name = theme["fonts"]["title"]
        p.font.color.rgb = RGBColor(*theme["colors"]["title"])
        p.alignment = PP_ALIGN.CENTER
        p.line_spacing = 1.2
        
        if slides_content and slides_content[0]:
            first_text = next(
                (s['content'] for s in slides_content[0] if s['type'] in ['h1', 'h2', 'text']),
                f"Generated in {mode} mode"
            )
            subtitle_box = title_slide.shapes.add_textbox(
                Inches(MARGIN_LEFT),
                Inches(SLIDE_HEIGHT * 0.5),
                Inches(CONTENT_WIDTH),
                Inches(1.0)
            )
            tf_sub = subtitle_box.text_frame
            tf_sub.word_wrap = True
            p_sub = tf_sub.paragraphs[0]
            p_sub.text = truncate_text(first_text, 80)
            p_sub.font.size = Pt(24)
            p_sub.font.name = theme["fonts"]["text"]
            p_sub.font.color.rgb = RGBColor(*theme["colors"]["subtitle"])
            p_sub.alignment = PP_ALIGN.CENTER
        
        # Content slides
        for slide_idx, slide_sections in enumerate(slides_content):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            set_background(slide, theme["colors"]["background"])
            add_decorative_shapes(slide, theme, "top")
            
            slide_title = "Overview"
            for section in slide_sections:
                if section['type'] in ['h1', 'h2']:
                    slide_title = section['content']
                    break
            
            print(f"📄 Slide {slide_idx + 1}: {truncate_text(slide_title, 60)}")
            
            # Title background
            title_bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(MARGIN_LEFT),
                Inches(MARGIN_TOP),
                Inches(CONTENT_WIDTH),
                Inches(TITLE_AREA_HEIGHT + 0.1)
            )
            title_bg.fill.solid()
            title_bg.fill.fore_color.rgb = RGBColor(*theme["colors"]["accent_light"])
            title_bg.line.fill.background()
            
            # Title text
            title_box = slide.shapes.add_textbox(
                Inches(MARGIN_LEFT + 0.2),
                Inches(MARGIN_TOP + 0.05),
                Inches(CONTENT_WIDTH - 0.4),
                Inches(TITLE_AREA_HEIGHT)
            )
            tf_title = title_box.text_frame
            tf_title.word_wrap = True
            p_title = tf_title.paragraphs[0]
            p_title.text = truncate_text(slide_title, 80)
            p_title.font.size = Pt(36)
            p_title.font.bold = True
            p_title.font.name = theme["fonts"]["heading"]
            p_title.font.color.rgb = RGBColor(*theme["colors"]["slide_title"])
            p_title.alignment = PP_ALIGN.LEFT
            
            # Image handling
            image_added = False
            if include_images and slide_idx > 0:
                print(f"  🖼️  Fetching image...")
                img_url = get_image_for_slide(slide_title)
                if img_url:
                    try:
                        img_data = download_image(img_url)
                        if img_data:
                            img_data.seek(0)
                            calc_width, calc_height = 3.2, 4.2
                            img_left = Inches(MARGIN_LEFT + IMAGE_START_LEFT + IMAGE_GAP)
                            img_top = Inches(CONTENT_AREA_TOP)
                            
                            slide.shapes.add_picture(
                                img_data,
                                img_left,
                                img_top,
                                width=Inches(calc_width),
                                height=Inches(calc_height)
                            )
                            image_added = True
                            print(f"  ✅ Image added!")
                    except Exception as e:
                        print(f"  ❌ Error adding image: {e}")
            
            # Font sizes
            font_config = calculate_dynamic_font_sizes(slide_sections, mode, image_added)
            if customizations.get('font_size'):
                font_size = int(customizations['font_size'])
                font_config['body'] = font_size
                font_config['h3'] = font_size + 2
            
            # Content area
            content_width = TEXT_WIDTH_WITH_IMAGE if image_added else TEXT_WIDTH_WITHOUT_IMAGE
            content_left = MARGIN_LEFT
            content_top = CONTENT_AREA_TOP
            content_height = AVAILABLE_CONTENT_HEIGHT
            
            content_box = slide.shapes.add_textbox(
                Inches(content_left),
                Inches(content_top),
                Inches(content_width),
                Inches(content_height)
            )
            tf_content = content_box.text_frame
            tf_content.word_wrap = True
            tf_content.vertical_anchor = MSO_ANCHOR.TOP
            tf_content.auto_size = MSO_AUTO_SIZE.NONE
            
            # Add content
            first_para = True
            for section in slide_sections:
                if section['type'] == 'h1' or section['type'] == 'h2':
                    continue
                
                if not first_para:
                    tf_content.add_paragraph()
                
                if section['type'] == 'h3':
                    p = tf_content.paragraphs[-1] if not first_para else tf_content.paragraphs[0]
                    p.text = truncate_text(section['content'], 250)
                    p.font.size = Pt(font_config['h3'])
                    p.font.bold = True
                    p.font.name = theme["fonts"]["heading"]
                    p.font.color.rgb = RGBColor(*theme["colors"]["h3"])
                    p.space_after = Pt(font_config['spacing'])
                    p.line_spacing = font_config['line_spacing']
                
                elif section['type'] == 'bullet':
                    p = tf_content.paragraphs[-1] if not first_para else tf_content.paragraphs[0]
                    p.text = truncate_text(section['content'], 250)
                    p.level = min(section['level'], 2)
                    p.font.size = Pt(font_config['body'])
                    p.font.name = theme["fonts"]["text"]
                    p.font.color.rgb = RGBColor(*theme["colors"]["bullet"])
                    p.space_after = Pt(font_config['spacing'])
                    p.line_spacing = font_config['line_spacing']
                
                elif section['type'] == 'text':
                    p = tf_content.paragraphs[-1] if not first_para else tf_content.paragraphs[0]
                    p.text = truncate_text(section['content'], 250)
                    p.font.size = Pt(font_config['body'])
                    p.font.name = theme["fonts"]["text"]
                    p.font.color.rgb = RGBColor(*theme["colors"]["text"])
                    p.space_after = Pt(font_config['spacing'] + 2)
                    p.line_spacing = font_config['line_spacing']
                
                first_para = False
        
        # Save presentation to disk (for direct use in video generation)
        import time
        timestamp = str(int(time.time()))
        final_filename = f"{timestamp}_{filename}.pptx"
        ppt_path = os.path.join(PPT_UPLOAD_DIR, final_filename)
        prs.save(ppt_path)
        
        # Determine the relative path for the frontend
        # We also want to return the blob so they can download it
        with open(ppt_path, 'rb') as f:
            ppt_bytes_data = f.read()
            
        print(f"\n✅ Presentation created successfully!")
        print(f"📊 Total slides: {len(prs.slides)}")
        print(f"📁 Saved to: {ppt_path}")
        print(f"📦 File size: {len(ppt_bytes_data) / 1024:.1f} KB\n")
        
        # We'll return JSON if the request headers say so, or just return the blob if they want attachment
        # Actually, to make it work with both, we'll return a JSON containing the path AND the download URL
        # But wait, send_file is more convenient for immediate download.
        # How about we return the path in custom headers? Or just return JSON?
        
        # PPT generation complete. immediate file delivery follows.

        # Let's check if the client expects JSON
        if request.headers.get('Accept') == 'application/json':
            return jsonify({
                "success": True,
                "ppt_path": ppt_path,
                "filename": final_filename,
                "message": "PPT generated successfully on server"
            })
        
        # Otherwise return the file (default for browser downloads)
        response = send_file(
            io.BytesIO(ppt_bytes_data),
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            as_attachment=True,
            download_name=f"{filename}.pptx"
        )
        response.headers['X-PPT-Path'] = ppt_path
        response.headers['X-PPT-Filename'] = final_filename
        response.headers['Access-Control-Expose-Headers'] = 'X-PPT-Path, X-PPT-Filename'
        return response
    
    except Exception as e:
        print(f"❌ Error generating PowerPoint: {str(e)}")
        import traceback
        traceback.print_exc()
        return jsonify({"error": str(e)}), 500

# ============ API ROUTES - VIDEO GENERATION ============
@app.route("/api/voices", methods=["GET"])
def get_voices():
    """Get available voice options"""
    if not VIDEO_GENERATION_ENABLED or not video_pipeline:
        return jsonify({
            "success": False,
            "error": "Video generation is not enabled"
        }), 503
    
    try:
        voices = video_pipeline.get_available_voices()
        return jsonify({
            "success": True,
            "voices": voices
        })
    except Exception as e:
        return jsonify({
            "success": False,
            "error": str(e)
        }), 500

@app.route("/api/upload-files", methods=["POST"])
def upload_files():
    """Upload PPT, face image, and optional audio for video generation"""
    if not VIDEO_GENERATION_ENABLED:
        return jsonify({
            "success": False,
            "error": "Video generation is not enabled"
        }), 503
    
    try:
        ppt_file = request.files.get('ppt')
        face_file = request.files.get('face')
        audio_file = request.files.get('audio')
        
        # Determine based on request.files if either is present
        if not ppt_file and not face_file and not audio_file:
            return jsonify({
                "success": False,
                "error": "At least one file (PPT, Face, or Audio) is required"
            }), 400
        
        # Use existing directories or determine which one to use
        # PPTs go to UPLOAD_FOLDER (or specifically PPT_UPLOAD_DIR if it exists)
        # Faces go to FACE_UPLOAD_DIR
        
        import time
        timestamp = str(int(time.time()))
        response_data = {"success": True}
        
        if ppt_file:
            # We assume UPLOAD_FOLDER or PPT_UPLOAD_DIR
            target_dir = globals().get('PPT_UPLOAD_DIR', UPLOAD_FOLDER)
            filename = f"{timestamp}_{secure_filename(ppt_file.filename)}"
            path = os.path.join(target_dir, filename)
            ppt_file.save(path)
            response_data["ppt_path"] = path

            # Save to DB if userId provided
            user_id = request.form.get('userId')
            if user_id:
                try:
                    db = get_db()
                    if db:
                        db.ppts.insert_one({
                            'userId': user_id,
                            'filename': ppt_file.filename,
                            'path': path,
                            'title': ppt_file.filename,
                            'createdAt': time.time(),
                            'type': 'uploaded'
                        })
                except Exception as dbe:
                    print(f"⚠️ Could not save uploaded PPT to DB: {dbe}")
            
        if face_file:
            target_dir = globals().get('FACE_UPLOAD_DIR', UPLOAD_FOLDER)
            filename = f"{timestamp}_{secure_filename(face_file.filename)}"
            path = os.path.join(target_dir, filename)
            face_file.save(path)
            response_data["face_path"] = path

        if audio_file:
            # Create audio upload directory if needed
            audio_upload_dir = os.path.join(UPLOAD_FOLDER, 'audio')
            os.makedirs(audio_upload_dir, exist_ok=True)
            filename = f"{timestamp}_{secure_filename(audio_file.filename)}"
            path = os.path.join(audio_upload_dir, filename)
            audio_file.save(path)
            response_data["audio_path"] = path
            
        return jsonify(response_data)
    
    except Exception as e:
        print(f"❌ Upload error: {e}")
        return jsonify({
            "success": False,
            "error": str(e)
        }), 500

@app.route("/api/generate-video", methods=["POST"])
def generate_video():
    """Generate video from PPT and face image"""
    if not VIDEO_GENERATION_ENABLED or not video_pipeline:
        return jsonify({
            "success": False,
            "error": "Video generation is not enabled. Please install required modules."
        }), 503
    
    try:
        data = request.json
        
        ppt_path = data.get('ppt_path')
        face_path = data.get('face_path')
        audio_path = data.get('audio_path')
        u_id = data.get('userId')

        # Generate JOB ID early for tracking
        job_id = f"{int(time.time())}_{os.path.splitext(os.path.basename(ppt_path))[0]}"
        
        # PERSIST active job to user profile if signed in
        if u_id:
            try:
                db_conn = get_db()
                if db_conn is not None:
                    db_conn.users.update_one({'clerkId': u_id}, {'$set': {'lastActiveJob': job_id}})
                    print(f"🔗 Linked Job {job_id} to User {u_id}")
            except Exception as dbe:
                print(f"⚠️ Failed to link job to user: {dbe}")
        
        if not ppt_path or not face_path:
            return jsonify({
                "success": False,
                "error": "PPT path and face path are required"
            }), 400
        
        # Check if files exist
        if not os.path.exists(ppt_path):
            return jsonify({
                "success": False,
                "error": f"PPT file not found: {ppt_path}"
            }), 404
        
        if not os.path.exists(face_path):
            return jsonify({
                "success": False,
                "error": f"Face image not found: {face_path}"
            }), 404

        if audio_path and not os.path.exists(audio_path):
            return jsonify({
                "success": False,
                "error": f"Custom audio file not found: {audio_path}"
            }), 404
        
        # Get options
        options = {
            'voice_id': data.get('voice_id', 0),
            'slang_level': data.get('slang_level', 'medium'),
            'quality': data.get('quality', 'medium'),
            'tts_engine': data.get('tts_engine', 'edge'),
            'audio_path': audio_path
        }
        
        # Function to run in background
        def run_pipeline():
            try:
                # Process the video
                video_pipeline.process(ppt_path, face_path, options, job_id=job_id)
                
                # PERSIST TO DATABASE upon completion
                if u_id:
                    try:
                        db = get_db()
                        if db is not None:
                            # Construct video data for gallery
                            video_title = os.path.splitext(os.path.basename(ppt_path))[0]
                            video_url = f"/api/download-video/{job_id}/final"
                            
                            session_data = {
                                "videoUrl": video_url,
                                "scriptUrl": f"/api/download-video/{job_id}/script",
                                "audioUrl": f"/api/download-video/{job_id}/audio",
                                "summary_url": f"/api/download-video/{job_id}/summary",
                                "jobId": job_id,
                                "userId": u_id,
                                "title": video_title
                            }
                            
                            video_doc = {
                                'userId': u_id,
                                'videoId': job_id,
                                'videoUrl': video_url,
                                'videoData': json.dumps(session_data),
                                'title': video_title,
                                'createdAt': time.time()
                            }
                            
                            # Check if already exists to avoid duplicates
                            if not db.videos.find_one({'videoId': job_id}):
                                db.videos.insert_one(video_doc)
                                print(f"✅ Video {job_id} auto-saved to User {u_id}'s gallery from backend")
                                
                                # Also clear the active job marker from the user profile
                                db.users.update_one({'clerkId': u_id}, {'$unset': {'lastActiveJob': ""}})
                    except Exception as db_err:
                        print(f"⚠️ Backend auto-save failed: {db_err}")

            except Exception as e:
                print(f"❌ BACKGROUND PIPELINE ERROR: {e}")
                # Save error to file so status poller can see it
                try:
                    with open(os.path.join(f'uploads/outputs/{job_id}', 'error.txt'), 'w') as f:
                        f.write(str(e))
                except: pass

        # Start processing in thread
        thread = threading.Thread(target=run_pipeline)
        thread.start()

        return jsonify({
            "success": True,
            "status": "processing",
            "job_id": job_id,
            "message": "Video generation started in background"
        })
    
    except Exception as e:
        print(f"❌ Video generation error: {str(e)}")
        import traceback
        traceback.print_exc()
        return jsonify({
            "success": False,
            "error": str(e),
            "traceback": traceback.format_exc()
        }), 500

@app.route("/api/download-video/<job_id>/<file_type>", methods=["GET"])
def download_video_file(job_id, file_type):
    """Download generated video files"""
    try:
        base_path = f'uploads/outputs/{job_id}'
        
        if not os.path.exists(base_path):
            return jsonify({"error": "Job not found"}), 404
        
        if file_type == 'final':
            file_path = os.path.join(base_path, 'final_lesson.mp4')
            mimetype = 'video/mp4'
            download_name = f'{job_id}_video.mp4'
        elif file_type == 'script':
            file_path = os.path.join(base_path, 'script.txt')
            mimetype = 'text/plain'
            download_name = f'{job_id}_script.txt'
        elif file_type == 'audio':
            file_path = os.path.join(base_path, 'narration.wav')
            
            # --- STRICT AUDIO VALIDATION & RE-ENCODING (Pydub + FFmpeg) ---
            print(f"\n🎧 AUDIO GENERATED")
            print(f"- Path: {file_path}")
            exists = os.path.exists(file_path)
            print(f"- Exists: {exists}")
            
            if not exists:
                raise RuntimeError("Audio generation failed: file does NOT exist")
            
            size = os.path.getsize(file_path)
            print(f"- File size: {size} bytes")
            if size < 5000:
                 raise RuntimeError("Audio generation failed: file size too small")

            print(f"\n🔄 AUDIO RE-ENCODING")
            print("- Target: PCM 16-bit WAV (44.1kHz)")
            
            from pydub import AudioSegment
            try:
                audio = AudioSegment.from_file(file_path)
                # Force conversion to PCM 16-bit, 44.1kHz, Mono (standard for educators)
                audio = audio.set_frame_rate(44100).set_sample_width(2).set_channels(1)
                audio.export(file_path, format="wav")
            except Exception as conv_err:
                raise RuntimeError(f"Corrupted WAV after conversion: {conv_err}")

            print(f"\n🧪 AUDIO VALIDATION")
            new_size = os.path.getsize(file_path)
            duration = len(audio) / 1000.0
            print(f"- File path: {file_path}")
            print(f"- Size: {new_size} bytes")
            print(f"- Duration: {duration:.2f}s")
            
            print(f"\n📦 DOWNLOAD READY")
            print(f"- Path: {file_path}")
            print(f"- MIME type: audio/wav")
            print(f"- Size: {new_size}")
            
            print("\n✅ AUDIO FILE READY FOR DOWNLOAD")
            print("- Fully playable WAV file")
            # -------------------------------------------------------------

            mimetype = 'audio/wav'
            download_name = f'eduface_{job_id}_narration.wav'
        elif file_type == 'summary':
            file_path = os.path.join(base_path, 'summary.txt')
            mimetype = 'text/plain'
            download_name = f'{job_id}_summary.txt'
        else:
            return jsonify({"error": "Invalid file type"}), 400
        
        if not os.path.exists(file_path):
            return jsonify({"error": f"File not found: {file_type}"}), 404
        
        return send_file(
            file_path,
            mimetype=mimetype,
            as_attachment=(file_type != 'final'), # Video must stream inline, not download
            download_name=download_name
        )
    
    except Exception as e:
        print(f"❌ Download error: {e}")
        return jsonify({"error": str(e)}), 500

@app.route("/api/video-status/<job_id>", methods=["GET"])
def get_video_status(job_id):
    """Check the status of a video generation job"""
    try:
        base_path = f'uploads/outputs/{job_id}'
        if not os.path.exists(base_path):
            return jsonify({"status": "not_found"}), 404
        
        # Check for error file
        if os.path.exists(os.path.join(base_path, 'error.txt')):
            with open(os.path.join(base_path, 'error.txt'), 'r') as f:
                error_msg = f.read()
            return jsonify({
                "status": "error",
                "error": error_msg
            })
        
        # Determine progress based on file existence
        progress = 0
        step = "Initializing"
        
        if os.path.exists(os.path.join(base_path, 'slides')):
            progress = 15
            step = "Extracting PPT"
        
        if os.path.exists(os.path.join(base_path, 'script.txt')):
            progress = 35
            step = "Generated Script"
            
        if any(f.startswith('audio_') for f in os.listdir(base_path) if os.path.isfile(os.path.join(base_path, f))):
            progress = 50
            step = "Generating Audio"
            
        if any(f.startswith('lipsync_') for f in os.listdir(base_path) if os.path.isfile(os.path.join(base_path, f))):
            progress = 75
            step = "Syncing Animation"
            
        if os.path.exists(os.path.join(base_path, 'final_lesson.mp4')):
            progress = 100
            step = "Completed"
            return jsonify({
                "status": "completed",
                "progress": 100,
                "step": step,
                "video_url": f"/api/download-video/{job_id}/final",
                "script_url": f"/api/download-video/{job_id}/script",
                "audio_url": f"/api/download-video/{job_id}/audio",
                "summary_url": f"/api/download-video/{job_id}/summary",
                "job_id": job_id
            })
            
        return jsonify({
            "status": "processing",
            "progress": progress,
            "step": step
        })
        
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/api/active-job/<user_id>", methods=["GET"])
def get_user_active_job(user_id):
    """Retrieve the last active job for a user"""
    try:
        db = get_db()
        if db is None:
            return jsonify({"error": "Database not connected"}), 500
            
        user = db.users.find_one({'clerkId': user_id})
        if not user or 'lastActiveJob' not in user:
            return jsonify({"jobId": None})
            
        return jsonify({"jobId": user['lastActiveJob']})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

# ============ EDUFACE AI CHATBOT ============
@app.route("/api/chat", methods=["POST"])
def tutor_chat():
    '''Elite Groq-powered Eduface AI (Llama 3.3)'''
    try:
        data = request.json
        messages = data.get('messages', [])
        api_key = os.getenv("GROQ_API_KEY", "").strip()
        if not api_key:
            return jsonify({"success": False, "error": "GROQ_API_KEY not found"}), 500
            
        client = Groq(api_key=api_key)

        # 1. CLEAN CONTEXT (Metadata Scrubbing)
        if messages and messages[0]['role'] == 'system':
            msg_content = messages[0]['content']
            clean_c = re.sub(r'\[\d{2}:\d{2}\]|\[SCENE START\]', '', msg_content)
            clean_c = re.sub(r'(FACE|EYES|HEAD|HANDS|BODY|TIMING):.*?(?=(FACE|EYES|HEAD|HANDS|BODY|TIMING|TEXT:|\Z))', '', clean_c, flags=re.DOTALL)
            clean_c = re.sub(r'TEXT:\s*"([^"]*)"', r'\1', clean_c)
            clean_c = re.sub(r'\s+', ' ', clean_c).strip()
            messages[0]['content'] = f"Context: {clean_c}"

        # 2. ELITE NOTEBOOKLM-STYLE IDENTITY
        system_i = {
            "role": "system",
            "content": (
                "You are 'Eduface AI', an intelligent real-time learner assistant. "
                "YOUR CORE RULE: ALWAYS respond in professional MARKDOWN. "
                "\n\nBEHAVIOR:"
                "\n1. Explain like an elite teacher. Be clear, structured, and authoritative."
                "\n2. Use **bolding** for key terms and concepts."
                "\n3. Use bullet points or numbered lists for steps and examples."
                "\n4. Use ### Headers for section breaks."
                "\n5. Ensure your response is professional and readable with appropriate spacing."
                "\n6. Engagement: Always offer follow-up help."
            )
        }
        messages.insert(0, system_i)

        comp = client.chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=messages,
            temperature=0.7
        )
        return jsonify({"success": True, "message": {"role": "assistant", "content": comp.choices[0].message.content}})
    except Exception as e:
        import traceback
        traceback.print_exc()
        # SCRUBBED FALLBACK POLICY
        friendly_fallback = "I’m currently unable to access the full model response, but I can still help you understand the topic."
        return jsonify({
            "success": True, 
            "message": {
                "role": "assistant", 
                "content": f"{friendly_fallback}\n\nThis lesson covers specialized materials. Would you like me to explain a related key concept while I refresh my connection?"
            }
        }), 200

# ============ EDUFACE QUIZ SYSTEM ============
@app.route("/api/quiz", methods=["POST"])
def manage_quiz():
    '''Generates and evaluates quizzes using Groq'''
    try:
        data = request.json
        lesson_content = data.get('lesson_content', '')
        num_questions = data.get('num_questions', 5)
        difficulty = data.get('difficulty', 'medium')
        user_answers = data.get('user_answers', {})
        full_quiz = data.get('full_quiz', [])
        
        api_key = os.getenv("GROQ_API_KEY", "").strip()
        if not api_key:
            return jsonify({"success": False, "error": "GROQ_API_KEY not found"}), 500
            
        client = Groq(api_key=api_key)

        prompt = f"""
        ACT AS: 'Eduface Smart Assistant'
        
        INPUTS:
        - Lesson Content: {lesson_content}
        - Number of Questions: {num_questions}
        - Difficulty: {difficulty}
        - User's Provided Answers: {json.dumps(user_answers)}

        MODE: {'EVALUATION' if user_answers else 'GENERATION'}
        QUIZ_POOL_SIZE: {len(full_quiz) if full_quiz else num_questions}
        QUIZ_CONTENT: {json.dumps(full_quiz) if full_quiz else 'N/A'}

        TASK (GENERATION):
        If this is GENERATION mode, identify the core concepts from the Lesson Content.
        Create EXACTLY {num_questions} high-quality MCQs at {difficulty} level.
        Each question must have 4 options (A, B, C, D) and ONE correct answer.
        Tag each question with a specific 'concept' found in the text.
        Provide a meaningful 'explanation' for the correct answer.

        TASK (EVALUATION):
        If this is EVALUATION mode, grade the 'User's Provided Answers' against the 'QUIZ_CONTENT' (if provided).
        If 'QUIZ_CONTENT' is NOT provided, it's GENERATION mode.
        CALCULATION RULE:
        1. Correct: Number of answers in 'User's Provided Answers' that match correct_answer in 'QUIZ_CONTENT'.
        2. Total: MUST be {len(full_quiz) if full_quiz else num_questions}.
        3. Wrong: Total - Correct (EVERY question NOT in 'User's Provided Answers' is WRONG).
        4. Accuracy: (Correct / Total) * 100.
        
        ADVANCED AI LEARNING ANALYST LAYER (Crucial):
        Analyze the quiz data and identify deep patterns in HOW the student learns and makes mistakes:
        - Conceptual Understanding: Which core topics are definitively weak or strong.
        - Error Patterns: Track repeated mistakes in similar questions. Differentiate fundamental misunderstandings vs careless errors.
        - Thinking Behavior: Is the user rushing? Overthinking? Guessing based on distractor choices?
        
        Perform concept-wise analysis for ALL unique concepts to determine accuracy per concept.

        OUTPUT FORMAT: STRICT JSON ONLY. NO MARKDOWN.
        {{
          "quiz": [
            {{
              "id": 1,
              "question": "...",
              "options": {{"A": "...", "B": "...", "C": "...", "D": "..."}},
              "correct_answer": "A",
              "concept": "...",
              "difficulty": "{difficulty}",
              "explanation": "..."
            }}
          ],
          "evaluation": {{
            "score": {{"total": 0, "correct": 0, "wrong": 0, "accuracy": "0%"}},
            "concept_analysis": [
              {{"concept": "...", "accuracy": "0%", "level": "Strong/Moderate/Weak"}}
            ],
            "learning_insights": [
              "Insight 1 about thinking behaviour",
              "Insight 2 about error patterns",
              "Insight 3 about conceptual depth"
            ],
            "strengths": ["Strength 1", "Strength 2"],
            "weaknesses": ["Weakness 1", "Weakness 2"],
            "suggestions": ["Behavioral advice", "Specific practice needed"],
            "learning_level": "Beginner/Intermediate/Advanced"
          }}
        }}
        """

        # Groq with Llama 3.3 70B
        completion = client.chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=[{"role": "system", "content": "You are a professional educational assessment engine. Respond ONLY in valid JSON."},
                      {"role": "user", "content": prompt}],
            temperature=0.1, # Keep it deterministic for JSON
            response_format={"type": "json_object"}
        )
        
        quiz_data = json.loads(completion.choices[0].message.content)
        return jsonify({"success": True, "data": quiz_data})
    except Exception as e:
        import traceback
        traceback.print_exc()
        return jsonify({"success": False, "error": str(e)}), 500

@app.route("/api/quiz/evaluate", methods=["POST"])
def evaluate_quiz_detailed():
    """Generates high-quality per-question reviews for a completed quiz"""
    try:
        data = request.json
        questions = data.get('questions', [])
        
        if not questions:
            return jsonify({"success": False, "error": "Missing quiz data"}), 400
            
        api_key = os.getenv("GROQ_API_KEY", "").strip()
        if not api_key:
            return jsonify({"success": False, "error": "GROQ_API_KEY not found"}), 500
            
        client = Groq(api_key=api_key)

        prompt = f"""
        You are an expert AI tutor and learning coach.
        Analyze the following student quiz results and provide a structured, educational review for each question.

        QUIZ DATA (JSON):
        {json.dumps(questions)}

        INSTRUCTIONS:
        For EACH question:
        1. Determine correctness (Result: Correct/Incorrect).
        2. If CORRECT:
           - Start with: "Your answer is correct!"
           - Explain WHY the correct answer is right in extreme depth (200-300 words).
           - Provide advanced professional insights, historical context, or complex real-world use cases.
           - Ensure the explanation is highly elaborated and teaches the concept thoroughly.
        3. If INCORRECT:
           - Start with: "Your answer is incorrect."
           - State the correct answer clearly.
           - Explain WHY the correct answer is right and critically analyze why the student's answer was logically flawed (200-300 words).
           - Provide deep technical details to correct the misconception entirely.
        4. ADDITIONAL RESOURCES:
           - Provide an array of exactly 2 "additionalResources".
           - These must NOT be generic. They must be highly specialized, conceptual, and heavily differentiated from other questions.
           - Format strictly as: "Advanced Concept Name - Detailed 2-sentence description of exactly what technical details to study to master this specific niche edge-case."

        TONE: Encouraging, highly professional, deeply analytical, and academic SaaS-level.
        GOAL: Provide a premium, university-level technical breakdown equivalent to a senior engineer's review.

        OUTPUT FORMAT: STRICT JSON ONLY. NO MARKDOWN.
        {{
          "evaluations": [
            {{
              "question": "...",
              "userAnswer": "...",
              "correctAnswer": "...",
              "result": "...",
              "explanation": "...",
              "mistakeExplanation": "...",
              "additionalResources": ["...", "..."]
            }}
          ]
        }}
        """

        completion = client.chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=[
                {"role": "system", "content": "You are a professional AI tutor and learning coach. Respond ONLY in valid JSON."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.3,
            response_format={"type": "json_object"}
        )
        
        result = json.loads(completion.choices[0].message.content)
        return jsonify({"success": True, "evaluations": result.get("evaluations", [])})
        
    except Exception as e:
        print(f"❌ EVALUATION ERROR: {str(e)}")
        return jsonify({"success": False, "error": str(e)}), 500



@app.route("/api/export-notes", methods=["POST"])
def export_study_notes():
    """Generates and exports professional academic notes as PDF or Word"""
    try:
        data = request.json
        if not data:
            return jsonify({"success": False, "error": "No data provided"}), 400
            
        lesson_content = data.get('lesson_content', '')
        style = data.get('style', 'EXPLANATIVE')
        depth = data.get('depth', 'INTERMEDIATE')
        include_examples = data.get('includeExamples', True)
        include_key_points = data.get('includeKeyPoints', True)
        export_format = data.get('format', 'PDF').upper()

        if not lesson_content:
            return jsonify({"success": False, "error": "Missing lesson content"}), 400

        api_key = os.getenv("GROQ_API_KEY", "").strip()
        if not api_key:
            return jsonify({"success": False, "error": "GROQ_API_KEY not found"}), 500

        client = Groq(api_key=api_key)

        # 1. Generate the Note Content using the elite prompt
        prompt = f"""
        ACT AS: 'Elite academic writer, senior instructional designer, and domain expert'
        DOCUMENT TYPE: Premium Academic Study Guide
        
        INPUT SOURCE (Lesson Transcript/Content):
        {lesson_content}

        USER REQUIREMENTS:
        - STYLE: {style} (EXPLANATIVE = deep paragraphs, BULLET = structured hierarchy, HINTS = dense keywords)
        - DEPTH: {depth} (BASIC = conceptual clarity, INTERMEDIATE = comprehensive coverage, ADVANCED = technical/in-depth analysis)
        - INCLUDE EXAMPLES: {'YES' if include_examples else 'NO'}
        - INCLUDE KEY POINTS: {'YES' if include_key_points else 'NO'}

        🎯 CORE OBJECTIVE:
        Transform the input source into a HIGH-FIDELITY, TEXTBOOK-QUALITY study guide. 
        DO NOT provide a 'script' or 'summary'. Generate the ACTUAL EDUCATIONAL CONTENT.
        Use formal, academic language. Ensure every complex concept in the source is unpacked and explained clearly based on the requested depth.

        STRUCTURE SPECIFICATIONS:
        1. TITLE: High-level academic heading.
        2. OVERVIEW: A master-level synthesis of the topic.
        3. CHAPTERS: Breakdown of the content into logical, deep-dive academic chapters.
        4. EXAMPLES: (If requested) Concrete applications or case studies.
        5. CONCLUSION: Final knowledge consolidation.

        🚫 RESTRICTIONS:
        - NO meta-talk (e.g., 'Here is your notes')
        - NO markdown symbols
        - NO placeholders

        OUTPUT FORMAT: Valid JSON only.
        {{
          "title": "...",
          "full_text": "A single, massive string of perfectly formatted text suitable for a document. Use newlines (\\n) for spacing. Include all sections, headers, and bullet points as specified above."
        }}
        """

        completion = client.chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=[
                {"role": "system", "content": "You are a professional textbook author. Your goal is to produce deep, high-value educational documentation. Respond ONLY in valid JSON."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.4,
            response_format={"type": "json_object"}
        )
        
        result = json.loads(completion.choices[0].message.content)
        title = result.get('title', 'Academic Study Guide')
        final_text = result.get('full_text', 'Significant educational content could not be synthesized.')

        # 2. Format as Document
        output = io.BytesIO()
        mimetype = 'text/plain'
        filename = f"{title.replace(' ', '_')}.txt"
        
        if export_format == "PDF":
            try:
                # Optimized PDF handling with fpdf2/fpdf support
                try: from fpdf import FPDF
                except ImportError: from fpdf2 import FPDF
                
                class PDF(FPDF):
                    def header(self):
                        try: self.set_font('Arial', 'B', 15)
                        except: self.set_font('helvetica', 'B', 15)
                        self.cell(0, 10, title, 0, 1, 'C')
                        self.ln(8)
                
                pdf = PDF()
                pdf.set_auto_page_break(auto=True, margin=15)
                pdf.add_page()
                try: pdf.set_font('Arial', '', 11)
                except: pdf.set_font('helvetica', '', 11)
                
                # Pre-process text to avoid encoding errors while preserving structure
                pdf.multi_cell(0, 8, final_text.encode('latin-1', 'replace').decode('latin-1'))
                
                raw_out = pdf.output(dest='S')
                # Handle fpdf vs fpdf2 output return types
                if isinstance(raw_out, (bytearray, bytes)):
                    output.write(raw_out)
                else:
                    output.write(raw_out.encode('latin-1'))
                    
                mimetype, filename = 'application/pdf', f"Study_Notes_{title.replace(' ', '_')}.pdf"
                
            except Exception as e:
                # Try ReportLab as second option
                try:
                    from reportlab.lib.pagesizes import letter
                    from reportlab.pdfgen import canvas
                    c = canvas.Canvas(output, pagesize=letter)
                    c.setFont("Helvetica-Bold", 16)
                    c.drawCentredString(300, 750, title)
                    c.setFont("Helvetica", 10)
                    textobject = c.beginText(50, 720)
                    for line in final_text.split('\n'):
                        textobject.textLine(line)
                    c.drawText(textobject)
                    c.showPage()
                    c.save()
                    mimetype, filename = 'application/pdf', f"{title.replace(' ', '_')}.pdf"
                except:
                    print(f"⚠️ All PDF Libs Missing: {e}")
                    # Final fallback to beautifully formatted plain text
                    header = f"{'='*50}\n{title.upper()}\n{'='*50}\n\n"
                    output.write((header + final_text).encode('utf-8'))
                    mimetype, filename = 'text/plain', f"{title.replace(' ', '_')}.txt"
        
        else: # WORD (DOCX)
            try:
                try: from docx import Document; from docx.shared import Pt
                except ImportError: from python_docx import Document; from docx.shared import Pt
                
                doc = Document()
                # Set Title
                title_run = doc.add_heading(title, 0).runs[0]
                title_run.font.size = Pt(20)
                
                # Split content into sections
                lines = final_text.split('\n')
                for line in lines:
                    line = line.strip()
                    if not line: continue
                    
                    # Logic to find headers
                    if line.isupper() and len(line) < 80:
                        doc.add_heading(line, level=1)
                    elif line.startswith('SECTION') or line.startswith('CHAPTER'):
                        doc.add_heading(line, level=1)
                    else:
                        p = doc.add_paragraph(line)
                        p.paragraph_format.space_after = Pt(10)
                
                doc.save(output)
                mimetype, filename = 'application/vnd.openxmlformats-officedocument.wordprocessingml.document', f"Study_Notes_{title.replace(' ', '_')}.docx"
                
            except Exception as e:
                print(f"⚠️ Word Export Fallback Active: {e}")
                html_body = []
                for b in final_text.split('\n\n'):
                    if b.strip():
                        if (b.isupper() and len(b) < 80) or b.startswith('SECTION'):
                            html_body.append(f"<h2 style='color: #1a365d; border-bottom: 1px solid #e2e8f0; margin-top: 30px; font-size: 1.4em;'>{b.strip()}</h2>")
                        else:
                            html_body.append(f"<p style='line-height: 1.7; color: #334155; margin-bottom: 12px;'>{b.strip().replace('\\n', '<br/>')}</p>")
                            
                html_content = f"""
                <html xmlns:o='urn:schemas-microsoft-com:office:office' xmlns:w='urn:schemas-microsoft-com:office:word' xmlns='http://www.w3.org/TR/REC-html40'>
                    <head><meta charset='UTF-8'></head>
                    <body style='font-family: Calibri, sans-serif; padding: 1in; max-width: 800px; margin: auto;'>
                        <h1 style='text-align: center; color: #2563eb; font-size: 2.2em; margin-bottom: 0.5in;'>{title}</h1>
                        {''.join(html_body)}
                        <hr style='margin-top: 50px; border: 0; border-top: 1px solid #cbd5e1;'/>
                        <p style='text-align: center; color: #94a3b8; font-size: 0.8em;'>Generated by Eduface AI - Premium Study Guide</p>
                    </body>
                </html>
                """
                output.write(html_content.encode('utf-8'))
                mimetype, filename = 'application/msword', f"Study_Notes_{title.replace(' ', '_')}.doc"

        output.seek(0)
        
        return send_file(
            output,
            mimetype=mimetype,
            as_attachment=True,
            download_name=filename
        )
        
    except Exception as e:
        import traceback
        traceback.print_exc()
        return jsonify({"success": False, "error": str(e)}), 500

# ============ HEALTH CHECK ============
@app.route("/api/health", methods=["GET"])
def health_check():
    """Health check endpoint"""
    ollama_ok, models = check_ollama_connection()
    
    return jsonify({
        "status": "healthy" if ollama_ok else "degraded",
        "ollama_running": ollama_ok,
        "ollama_url": OLLAMA_BASE_URL,
        "ai_model": AI_MODEL,
        "available_models": models,
        "available_modes": list(MODE_PROMPTS.keys()),
        "video_generation_enabled": VIDEO_GENERATION_ENABLED,
        "image_apis": {
            "unsplash": bool(UNSPLASH_API_KEY),
            "pexels": bool(PEXELS_API_KEY),
            "pixabay": bool(PIXABAY_API_KEY)
        }
    })

# ==========================================
# Chat Sharing Endpoints
# ==========================================
@app.route('/api/share-chat', methods=['POST'])
def share_chat():
    try:
        data = request.json
        messages = data.get('messages', [])
        visibility = data.get('visibility', 'public')
        
        chat_id = str(uuid.uuid4())
        
        # Ensure only safe client data is stored
        safe_messages = []
        for msg in messages:
            safe_messages.append({
                'role': msg.get('role'),
                'content': msg.get('content')
            })
            
        chat_data = {
            'id': chat_id,
            'visibility': visibility,
            'messages': safe_messages,
            'timestamp': time.time(),
            'read_only': True
        }
        
        file_path = os.path.join(SHARED_CHATS_DIR, f"{chat_id}.json")
        with open(file_path, 'w', encoding='utf-8') as f:
            json.dump(chat_data, f)
            
        return jsonify({
            'success': True,
            'share_id': chat_id,
            'share_url': f"http://localhost:5173/shared/{chat_id}"
        })
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500

@app.route('/api/shared-chat/<chat_id>', methods=['GET'])
def get_shared_chat(chat_id):
    try:
        # Prevent directory traversal
        safe_id = secure_filename(f"{chat_id}.json")
        file_path = os.path.join(SHARED_CHATS_DIR, safe_id)
        
        if not os.path.exists(file_path):
            return jsonify({'success': False, 'error': 'Shared chat not found'}), 404
            
        with open(file_path, 'r', encoding='utf-8') as f:
            chat_data = json.load(f)
            
        return jsonify({
            'success': True,
            'chat': chat_data
        })
    except Exception as e:
        return jsonify({'success': False, 'error': 'Failed to load chat'}), 500

# ==========================================
# MongoDB User & Video Sync Endpoints
# ==========================================
try:
    from db import get_db
except ImportError:
    get_db = lambda: None

@app.route('/api/users', methods=['POST'])
def sync_user():
    try:
        data = request.json
        if not data or not data.get('clerkId'):
            return jsonify({'success': False, 'error': 'clerkId is required'}), 400
            
        db = get_db()
        if db is None:
            return jsonify({'success': False, 'error': 'Database not connected'}), 500
            
        clerk_id = data['clerkId']
        user_data = {
            'name': data.get('name', ''),
            'email': data.get('email', ''),
            'profileImage': data.get('profileImage', ''),
            'lastLogin': time.time()
        }
        
        # update or insert
        db.users.update_one(
            {'clerkId': clerk_id},
            {
                '$set': user_data,
                '$setOnInsert': {'createdAt': time.time()}
            },
            upsert=True
        )
        
        return jsonify({'success': True, 'message': 'User synced successfully'})
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500

@app.route('/api/videos', methods=['POST'])
def save_video():
    try:
        data = request.json
        print(f"📥 SAVE VIDEO REQUEST: {data}")
        db = get_db()
        if db is None:
            return jsonify({'success': False, 'error': 'Database not connected'}), 500
            
        u_id = data.get('userId')
        v_id = data.get('videoId')
        
        if not u_id or not v_id:
            print(f"❌ REJECTED: userId={u_id}, videoId={v_id}")
            return jsonify({'success': False, 'error': 'userId and videoId required'}), 400
            
        # check if video already exists
        existing = db.videos.find_one({'videoId': v_id})
        if existing:
            print(f"ℹ️ Video {v_id} already in gallery")
            return jsonify({'success': True, 'message': 'Video already exists'})
            
        video_data = {
            'userId': data['userId'],
            'videoId': data['videoId'],
            'videoUrl': data.get('videoUrl', ''),
            'videoData': data.get('videoData', ''),
            'title': data.get('title', 'Generated Video'),
            'createdAt': data.get('createdAt', time.time())
        }
        
        db.videos.insert_one(video_data)
        return jsonify({'success': True, 'message': 'Video saved to DB'})
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500

@app.route('/api/videos/<clerk_id>', methods=['GET'])
def get_user_videos(clerk_id):
    try:
        db = get_db()
        if db is None:
            return jsonify({'success': False, 'error': 'Database not connected'}), 500
            
        videos = list(db.videos.find({'userId': clerk_id}).sort('createdAt', -1))
        
        # Convert ObjectId to string for JSON serialization
        for v in videos:
            v['_id'] = str(v['_id'])
            
        return jsonify({
            'success': True,
            'videos': videos
        })
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500

from bson.objectid import ObjectId

@app.route('/api/videos/<video_id>', methods=['PUT'])
def rename_video(video_id):
    try:
        data = request.json
        new_title = data.get('title')
        if not new_title:
            return jsonify({'success': False, 'error': 'New title is required'}), 400
            
        db = get_db()
        if db is None:
            return jsonify({'success': False, 'error': 'Database not connected'}), 500
            
        result = db.videos.update_one(
            {'_id': ObjectId(video_id)},
            {'$set': {'title': new_title}}
        )
        
        if result.modified_count == 0:
            return jsonify({'success': False, 'error': 'Video not found or title unchanged'}), 404
            
        return jsonify({'success': True, 'message': 'Video renamed successfully'})
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500

@app.route('/api/videos/<video_id>', methods=['DELETE'])
def delete_video(video_id):
    try:
        db = get_db()
        if db is None:
            return jsonify({'success': False, 'error': 'Database not connected'}), 500
            
        # Optional: delete physical files from storage if needed. 
        # For now, just delete DB record to remove from gallery.
        result = db.videos.delete_one({'_id': ObjectId(video_id)})
        
        if result.deleted_count == 0:
            return jsonify({'success': False, 'error': 'Video not found'}), 404
            
        return jsonify({'success': True, 'message': 'Video deleted successfully'})
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500

@app.route('/api/ppts/<clerk_id>', methods=['GET'])
def get_user_ppts(clerk_id):
    try:
        db = get_db()
        if db is None:
            return jsonify({'success': False, 'error': 'Database not connected'}), 500
            
        ppts = list(db.ppts.find({'userId': clerk_id}).sort('createdAt', -1))
        
        # Convert ObjectId to string
        for p in ppts:
            p['_id'] = str(p['_id'])
            
        return jsonify({
            'success': True,
            'ppts': ppts
        })
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500

# ============ START SERVER ============
if __name__ == "__main__":
    app.run(debug=False, host="0.0.0.0", port=5000)