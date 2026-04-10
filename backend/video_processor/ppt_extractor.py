import os
import sys
import subprocess
import glob
import shutil
from pptx import Presentation

class PPTExtractor:
    def __init__(self, ppt_path):
        self.ppt_path = ppt_path
        self.presentation = Presentation(ppt_path)

    def extract_text(self):
        """Extract all text content from PPT slides including tables, groups, and paragraph-level detail"""
        slides_content = []
        for slide_num, slide in enumerate(self.presentation.slides, 1):
            slide_text = []

            def _process_shape(shape):
                # 1. Handle shapes with text frames (most text)
                if hasattr(shape, "text_frame") and shape.text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_text.append(text)
                
                # 2. Handle tables specifically
                elif shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            # Extract text from every paragraph in every cell
                            for paragraph in cell.text_frame.paragraphs:
                                text = paragraph.text.strip()
                                if text:
                                    slide_text.append(text)
                
                # 3. Handle groups (Recursive)
                elif shape.shape_type == 6: # Group shape
                    for subshape in shape.shapes:
                        _process_shape(subshape)
                
                # 4. Handle Graphic Frames (SmartArt, Charts, etc.)
                elif shape.has_chart:
                    try:
                        if shape.chart.has_title:
                            slide_text.append(shape.chart.chart_title.text_frame.text)
                    except: pass
                
                # 5. Fallback for other text-bearing shapes
                elif hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())

            # Sort shapes by top/left position to maintain logical reading order
            sorted_shapes = sorted(slide.shapes, key=lambda s: (s.top if hasattr(s, 'top') else 0, s.left if hasattr(s, 'left') else 0))
            
            for shape in sorted_shapes:
                _process_shape(shape)

            notes = ""
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                if notes_slide.notes_text_frame:
                    notes = notes_slide.notes_text_frame.text

            # Deduplicate while preserving order (some shapes might double-report text)
            seen = set()
            distinct_text = []
            for t in slide_text:
                if t not in seen:
                    distinct_text.append(t)
                    seen.add(t)

            slides_content.append({
                'slide_number': slide_num,
                'title': distinct_text[0] if distinct_text else f"Slide {slide_num}",
                'content': '\n'.join(distinct_text),
                'notes': notes.strip()
            })
        return slides_content

    def get_slide_count(self):
        return len(self.presentation.slides)

    def export_slides_as_images(self, output_dir, width=1920, height=1080):
        """
        Export each PPTX slide as a PNG image.
        Tries MS PowerPoint COM first (100% fidelity on Windows).
        Tries LibreOffice headless second.
        Falls back to PIL-based renderer.
        Returns a sorted list of PNG file paths.
        """
        os.makedirs(output_dir, exist_ok=True)
        image_paths = []
        abs_ppt = os.path.abspath(self.ppt_path)
        abs_out = os.path.abspath(output_dir)

        # ── Attempt 1: MS PowerPoint COM (Windows Only, 100% Fidelity) ──
        try:
            import sys
            if sys.platform == 'win32':
                import win32com.client
                import pythoncom
                print("  [PROCESS] Using MS PowerPoint to export slides (100% fidelity)...")
                pythoncom.CoInitialize()
                powerpoint = win32com.client.Dispatch("PowerPoint.Application")
                presentation = powerpoint.Presentations.Open(abs_ppt, ReadOnly=True, WithWindow=False)
                
                for i, slide in enumerate(presentation.Slides, 1):
                    img_path = os.path.join(abs_out, f"slide_{i-1:03d}.png")
                    slide.Export(img_path, "PNG", width, height)
                    if os.path.exists(img_path):
                        image_paths.append(img_path)
                
                presentation.Close()
                try: powerpoint.Quit()
                except: pass
                pythoncom.CoUninitialize()
                
                if image_paths:
                    print(f"  [SUCCESS] Exported {len(image_paths)} slides perfectly via PowerPoint")
                    return sorted(image_paths)
        except Exception as com_err:
            print(f"  [WARNING] PowerPoint COM export failed: {com_err}")

        # ── Attempt 2: LibreOffice headless ─────────────────────────────
        try:
            lo_candidates = [
                r'C:\Program Files\LibreOffice\program\soffice.exe',
                r'C:\Program Files (x86)\LibreOffice\program\soffice.exe',
                'soffice',
            ]
            lo_exe = None
            for candidate in lo_candidates:
                path = shutil.which(candidate)
                if path:
                    lo_exe = path
                    break

            if lo_exe:
                print(f"  [PROCESS] Using LibreOffice to export slides...")
                # Convert all arguments to string to avoid complex type inference issues (PathLike recursion)
                subprocess.run(
                    [str(lo_exe), '--headless', '--convert-to', 'png', '--outdir', str(abs_out), str(abs_ppt)],
                    capture_output=True, text=True, timeout=120
                )
                base = os.path.splitext(os.path.basename(abs_ppt))[0]
                exported = sorted(glob.glob(os.path.join(abs_out, f"{base}*.png")))
                if exported:
                    print(f"  [SUCCESS] Exported {len(exported)} slides via LibreOffice")
                    return exported
        except Exception as lo_err:
            print(f"  [WARNING] LibreOffice export failed: {lo_err}")

        # ── Attempt 3: PIL renderer (always works) ──────────────────────
        try:
            from PIL import Image, ImageDraw, ImageFont

            prs = self.presentation
            emu_per_px = 914400 / 96          # 96 DPI
            slide_w = int(prs.slide_width  / emu_per_px)
            slide_h = int(prs.slide_height / emu_per_px)
            scale   = min(width / slide_w, height / slide_h)
            out_w   = int(slide_w * scale)
            out_h   = int(slide_h * scale)

            print(f"  [PROCESS] Using PIL renderer for slides ({out_w}×{out_h})...")
            for i, slide in enumerate(prs.slides):
                # 1. Background colour extraction (Hierarchical)
                bg_color = (255, 255, 255) # default white
                try:
                    # Check slide background
                    if hasattr(slide, 'background') and hasattr(slide.background, 'fill') and slide.background.fill.type == 1:
                        c = slide.background.fill.fore_color.rgb
                        bg_color = (c[0], c[1], c[2]) if type(c) == tuple else (c.r, c.g, c.b)
                    # Check layout background
                    elif hasattr(slide.slide_layout, 'background') and hasattr(slide.slide_layout.background, 'fill') and slide.slide_layout.background.fill.type == 1:
                        c = slide.slide_layout.background.fill.fore_color.rgb
                        bg_color = (c[0], c[1], c[2]) if type(c) == tuple else (c.r, c.g, c.b)
                    # Check master background
                    elif hasattr(slide.slide_master, 'background') and hasattr(slide.slide_master.background, 'fill') and slide.slide_master.background.fill.type == 1:
                        c = slide.slide_master.background.fill.fore_color.rgb
                        bg_color = (c[0], c[1], c[2]) if type(c) == tuple else (c.r, c.g, c.b)
                except Exception as e:
                    pass

                # If the theme has a dark background but python-pptx can't find it, it will be white.
                img  = Image.new('RGB', (out_w, out_h), bg_color)
                draw = ImageDraw.Draw(img)

                # Determine if background is "dark" to auto-adjust text color if missing
                brightness = (bg_color[0] * 299 + bg_color[1] * 587 + bg_color[2] * 114) / 1000
                is_dark_bg = brightness < 128
                default_text_color = (245, 245, 245) if is_dark_bg else (20, 20, 20)

                for shape in slide.shapes:
                    if not hasattr(shape, "text") or not shape.text.strip():
                        continue
                    try:
                        x = int(shape.left / emu_per_px * scale) if shape.left else 50
                        y = int(shape.top  / emu_per_px * scale) if shape.top else 50
                        font_size = max(14, int(20 * scale))
                        
                        try:
                            # Try to extract font color from the first run
                            text_color = default_text_color
                            if hasattr(shape, "text_frame") and shape.text_frame.paragraphs:
                                p = shape.text_frame.paragraphs[0]
                                if p.runs and p.runs[0].font.color and p.runs[0].font.color.type == 1:
                                    c = p.runs[0].font.color.rgb
                                    text_color = (c[0], c[1], c[2]) if type(c) == tuple else (c.r, c.g, c.b)
                        except Exception:
                            text_color = default_text_color

                        try:
                            font = ImageFont.truetype("arial.ttf", font_size)
                        except Exception:
                            font = ImageFont.load_default()
                            
                        # Basic text wrapping
                        import textwrap
                        wrapped_text = "\n".join(textwrap.wrap(shape.text.strip(), width=int(80 * (1.0/scale))))
                        draw.text((x + 10, y + 10), wrapped_text,
                                  fill=text_color, font=font)

                    except Exception:
                        pass

                # Centre on black canvas (letterbox)
                canvas  = Image.new('RGB', (width, height), (0, 0, 0))
                paste_x = (width  - out_w) // 2
                paste_y = (height - out_h) // 2
                canvas.paste(img, (paste_x, paste_y))

                img_path = os.path.join(output_dir, f'slide_{i:03d}.png')
                canvas.save(img_path, 'PNG')
                image_paths.append(img_path)

            print(f"  [SUCCESS] Exported {len(image_paths)} slides via PIL")
            return image_paths

        except Exception as pil_err:
            print(f"  [WARNING] PIL renderer failed: {pil_err}")
            return []